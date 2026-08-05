"""
RAG pipeline — document ingestion, hybrid retrieval, answer generation.

Changes vs previous version:
  - In-memory _answer_cache dict replaced with Redis (persistent, shared across workers).
  - Cache gracefully degrades to no-op if Redis is unavailable.
  - Per-user answer cache invalidation still works (DEL by key pattern).
"""
import os
import json
import hashlib
import logging
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import FastEmbedEmbeddings
from langchain_groq import ChatGroq
from langchain_chroma import Chroma
from langchain.chains import RetrievalQA

logger = logging.getLogger(__name__)

GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# Use /data (Render persistent disk) if available, else local directory
_DATA_DIR = "/data" if os.path.isdir("/data") else os.path.dirname(os.path.abspath(__file__))
CHROMA_DIR = os.getenv("CHROMA_DIR", os.path.join(_DATA_DIR, "chroma_db"))

# ── Redis cache ────────────────────────────────────────────────────────────────
# Answers are cached in Redis with a 1-hour TTL.
# Falls back silently to no caching if Redis is unreachable.

_CACHE_TTL = int(os.getenv("ANSWER_CACHE_TTL", "3600"))  # seconds
_CACHE_PREFIX = "askmynotes:answer"
_redis = None


def _get_redis():
    """Lazy-initialise Redis connection. Returns None if unavailable."""
    global _redis
    if _redis is not None:
        return _redis
    redis_url = os.getenv("REDIS_URL", "redis://localhost:6379/0")
    try:
        import redis as _redis_lib
        client = _redis_lib.from_url(
            redis_url,
            decode_responses=True,
            socket_connect_timeout=2,
            socket_timeout=2,
        )
        client.ping()
        _redis = client
        logger.info("Redis cache connected at %s", redis_url.split("@")[-1])
    except Exception as exc:
        logger.warning("Redis unavailable — answer cache disabled (%s)", exc)
        _redis = None
    return _redis


def _cache_key(user_email: str, question: str) -> str:
    u = hashlib.md5(user_email.encode()).hexdigest()[:12]
    q = hashlib.md5(question.lower().strip().encode()).hexdigest()
    return f"{_CACHE_PREFIX}:{u}:{q}"


def _cache_get(user_email: str, question: str) -> dict | None:
    r = _get_redis()
    if r is None:
        return None
    try:
        raw = r.get(_cache_key(user_email, question))
        return json.loads(raw) if raw else None
    except Exception:
        return None


def _cache_set(user_email: str, question: str, value: dict) -> None:
    r = _get_redis()
    if r is None:
        return
    try:
        r.setex(_cache_key(user_email, question), _CACHE_TTL, json.dumps(value))
    except Exception:
        pass  # cache write failure is non-fatal


def _cache_invalidate_user(user_email: str) -> None:
    """Delete all cached answers for a user (called when they upload new content)."""
    r = _get_redis()
    if r is None:
        return
    try:
        u = hashlib.md5(user_email.encode()).hexdigest()[:12]
        pattern = f"{_CACHE_PREFIX}:{u}:*"
        cursor = 0
        while True:
            cursor, keys = r.scan(cursor, match=pattern, count=100)
            if keys:
                r.delete(*keys)
            if cursor == 0:
                break
        logger.debug("Invalidated Redis cache for user %s", user_email)
    except Exception:
        pass  # cache invalidation failure is non-fatal


# ── Lazy loaders (prevent Render startup timeouts) ─────────────────────────────
_embeddings = None
_reranker   = None


def _get_embeddings():
    global _embeddings
    if _embeddings is None:
        _embeddings = FastEmbedEmbeddings(model_name="BAAI/bge-small-en-v1.5")
    return _embeddings


def _get_reranker():
    """Lazy-load Flashrank cross-encoder re-ranker (tiny model, no API key needed)."""
    global _reranker
    if _reranker is None:
        from langchain_community.document_compressors.flashrank_rerank import FlashrankRerank
        _reranker = FlashrankRerank(top_n=5)
    return _reranker


# ── Per-user vector store cache ────────────────────────────────────────────────
# Maps user_email → Chroma instance so we don't recreate connections on every call
_user_stores: dict = {}


def _collection_name(user_email: str) -> str:
    """Generate a safe, unique ChromaDB collection name from a user's email."""
    return "user_" + hashlib.md5(user_email.encode()).hexdigest()[:16]


def get_user_vectorstore(user_email: str) -> Chroma:
    """Return the Chroma collection for this user, creating it if needed."""
    if user_email not in _user_stores:
        _user_stores[user_email] = Chroma(
            collection_name=_collection_name(user_email),
            embedding_function=_get_embeddings(),
            persist_directory=CHROMA_DIR,
        )
    return _user_stores[user_email]


def _hash(text: str) -> str:
    return hashlib.md5(text.lower().strip().encode()).hexdigest()


# ── Multi-format document loader ───────────────────────────────────────────────

def _load_pptx(file_path: str) -> list:
    """Extract text from a PowerPoint file slide-by-slide."""
    from pptx import Presentation
    from langchain_core.documents import Document

    prs = Presentation(file_path)
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"source": file_path, "page": i, "slide": i + 1},
            ))
    return docs


def load_document(file_path: str) -> list:
    """
    Load a document based on its file extension.
    Supported: .pdf, .docx, .pptx, .txt
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return PyPDFLoader(file_path).load()

    elif ext == ".txt":
        return TextLoader(file_path, encoding="utf-8").load()

    elif ext == ".docx":
        from langchain_community.document_loaders import Docx2txtLoader
        return Docx2txtLoader(file_path).load()

    elif ext == ".pptx":
        return _load_pptx(file_path)

    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf .docx .pptx .txt")


# ── Core pipeline functions ────────────────────────────────────────────────────

def _build_hybrid_retriever(user_email: str, k: int = 12):
    """
    Combine BM25 keyword search (40%) with dense MMR semantic search (60%).
    Falls back to MMR-only if the collection is empty.
    """
    from langchain_community.retrievers import BM25Retriever
    from langchain.retrievers import EnsembleRetriever
    from langchain_core.documents import Document

    vs = get_user_vectorstore(user_email)
    all_data = vs.get(include=["documents", "metadatas"])

    if not all_data["documents"]:
        return vs.as_retriever(search_type="mmr", search_kwargs={"k": k // 2, "fetch_k": k})

    docs = [
        Document(page_content=text, metadata=meta)
        for text, meta in zip(all_data["documents"], all_data["metadatas"])
    ]
    bm25  = BM25Retriever.from_documents(docs, k=k)
    dense = vs.as_retriever(search_type="mmr", search_kwargs={"k": k, "fetch_k": k * 2})
    return EnsembleRetriever(retrievers=[bm25, dense], weights=[0.4, 0.6])


def get_relevant_docs(user_email: str, question: str, k: int = 5) -> list:
    """
    Advanced RAG retrieval:
      1. Hybrid search  — BM25 (keyword) + dense MMR (semantic)
      2. Cross-encoder re-ranking — Flashrank picks the best k results
    Falls back to plain MMR if either step fails.
    """
    from langchain.retrievers import ContextualCompressionRetriever
    try:
        hybrid   = _build_hybrid_retriever(user_email, k=k * 2)
        reranker = _get_reranker()
        retriever = ContextualCompressionRetriever(
            base_compressor=reranker,
            base_retriever=hybrid,
        )
        docs = retriever.invoke(question)
        return docs if docs else _mmr_fallback(user_email, question, k)
    except Exception:
        return _mmr_fallback(user_email, question, k)


def _mmr_fallback(user_email: str, question: str, k: int) -> list:
    """Plain MMR fallback when advanced retrieval fails."""
    vs = get_user_vectorstore(user_email)
    return vs.as_retriever(
        search_type="mmr", search_kwargs={"k": k, "fetch_k": k * 2}
    ).invoke(question)


def list_indexed_files(user_email: str) -> list:
    """Return unique filenames indexed in this user's collection."""
    try:
        vs = get_user_vectorstore(user_email)
        results = vs.get(include=["metadatas"])
        files = set()
        for meta in results.get("metadatas", []):
            src = meta.get("source", "")
            if src:
                files.add(os.path.basename(src))
        return sorted(list(files))
    except Exception:
        return []


def clear_vectorstore(user_email: str):
    """Delete all documents from this user's vector collection."""
    global _user_stores
    # Invalidate Redis cache for this user
    _cache_invalidate_user(user_email)

    try:
        vs = get_user_vectorstore(user_email)
        vs.delete_collection()
        # Remove from local store cache so it gets recreated fresh on next access
        _user_stores.pop(user_email, None)
    except Exception:
        pass


def ingest_document(file_path: str, user_email: str) -> int:
    """
    Load any supported document, split into chunks, and store in the
    user's private ChromaDB collection. Returns the number of chunks indexed.
    Invalidates Redis answer cache so stale answers aren't served.
    """
    pages = load_document(file_path)

    # Larger overlap (100) reduces context loss at chunk boundaries
    splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    chunks = splitter.split_documents(pages)

    vs = get_user_vectorstore(user_email)
    vs.add_documents(chunks)

    # Invalidate this user's Redis answer cache when new content is added
    _cache_invalidate_user(user_email)

    return len(chunks)


# Keep old name as an alias so nothing outside this module needs to change yet
def ingest_pdf(file_path: str, user_email: str) -> int:
    return ingest_document(file_path, user_email)


def get_answer(question: str, user_email: str) -> dict:
    """Advanced RAG: hybrid retrieval + re-ranking + LLM answer with citations."""

    # 1. Check Redis cache first
    cached = _cache_get(user_email, question)
    if cached is not None:
        return {**cached, "cached": True}

    # 2. Advanced retrieval: hybrid search + re-ranking
    docs = get_relevant_docs(user_email, question, k=5)
    context = "\n\n".join(d.page_content for d in docs)

    # 3. Build source citations
    sources = []
    seen_pages = set()
    for doc in docs:
        page = doc.metadata.get("page", 0) + 1
        filename = os.path.basename(doc.metadata.get("source", "unknown"))
        if page not in seen_pages:
            seen_pages.add(page)
            sources.append({
                "page": page,
                "file": filename,
                "snippet": doc.page_content[:200].strip(),
            })

    # 4. LLM answer generation
    llm = ChatGroq(
        api_key=GROQ_API_KEY,
        model=os.getenv("GROQ_MODEL", "llama-3.1-8b-instant"),
        temperature=0,
    )
    prompt = (
        "You are a study assistant. Answer based ONLY on the provided context. "
        "If the answer is not in the context, say so.\n\n"
        f"Context:\n{context}\n\nQuestion: {question}\nAnswer:"
    )
    result = llm.invoke(prompt)

    answer = {
        "answer": result.content,
        "sources": sources,
        "cached": False,
    }

    # 5. Store in Redis (TTL = ANSWER_CACHE_TTL seconds, default 1 hour)
    _cache_set(user_email, question, answer)
    return answer


def get_top_chunks(user_email: str, k: int = 8) -> str:
    """
    Retrieve the most relevant chunks for this user's notes using hybrid search.
    Used by quiz.py and study-guide endpoint.
    """
    docs = get_relevant_docs(
        user_email,
        "key concepts, definitions, important facts and examples",
        k=k,
    )
    return "\n\n".join(doc.page_content for doc in docs)


def delete_file_from_vectorstore(user_email: str, filename: str) -> int:
    """Remove all chunks for a specific file from the user's vector collection.
    Returns number of chunks deleted."""
    vs = get_user_vectorstore(user_email)
    results = vs.get(include=["metadatas"])
    ids_to_delete = [
        doc_id for doc_id, meta in zip(results["ids"], results["metadatas"])
        if os.path.basename(meta.get("source", "")) == filename
    ]
    if ids_to_delete:
        vs.delete(ids=ids_to_delete)
    # Invalidate Redis answer cache for this user
    _cache_invalidate_user(user_email)
    return len(ids_to_delete)
